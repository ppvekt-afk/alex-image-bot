import os
import re
import shutil
import asyncio
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

class PPTSkill:
    def __init__(self):
        self.output_dir = Path("generated_pptx")
        self.output_dir.mkdir(exist_ok=True)
        self.html_dir = Path("generated_html_ppts")
        self.html_dir.mkdir(exist_ok=True)
        
        possible_paths = [
            Path("open-design/design-templates/html-ppt"),
            Path("html-ppt-skill"),
        ]
        self.skill_path = None
        for p in possible_paths:
            if p.exists():
                self.skill_path = p
                break
        
        try:
            from apograph import convert_html_to_pptx
            self.apograph_available = True
        except ImportError:
            self.apograph_available = False
        
        if self.skill_path:
            print(f"✅ Навык презентаций загружен из: {self.skill_path}")
    
    def list_themes(self):
        if not self.skill_path:
            return ["minimal-white", "cyberpunk-neon", "corporate-clean"]
        themes_path = self.skill_path / "assets" / "themes"
        if themes_path.exists():
            return [f.stem for f in themes_path.glob("*.css")]
        return ["minimal-white", "cyberpunk-neon", "corporate-clean"]
    
    def list_animations(self):
        if not self.skill_path:
            return {"css": ["fade", "slide", "zoom"], "canvas_fx": []}
        anims_path = self.skill_path / "assets" / "animations"
        css_animations = []
        if (anims_path / "animations.css").exists():
            css_content = (anims_path / "animations.css").read_text(encoding='utf-8', errors='ignore')
            css_animations = re.findall(r'\.([\w-]+)-animation\s*\{', css_content)
        fx_animations = []
        if (anims_path / "fx").exists():
            fx_animations = [f.stem for f in (anims_path / "fx").glob("*.js")]
        return {"css": css_animations[:27], "canvas_fx": fx_animations[:20]}
    
    def create_html_ppt(self, topic, slides=5, theme="cyberpunk-neon", animation="fade"):
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_folder = self.html_dir / f"ppt_{timestamp}"
        output_folder.mkdir(exist_ok=True)
        
        if not self.skill_path:
            return None, None, "Навык презентаций не установлен"
        
        template_path = self.skill_path / "templates" / "full-decks" / "pitch-deck"
        if not template_path.exists():
            template_path = self.skill_path / "templates" / "deck.html"
        
        if template_path.exists():
            if template_path.is_dir():
                shutil.copytree(template_path, output_folder, dirs_exist_ok=True)
                index_path = output_folder / "index.html"
            else:
                shutil.copy(template_path, output_folder / "index.html")
                index_path = output_folder / "index.html"
            
            if index_path.exists():
                content = index_path.read_text(encoding='utf-8', errors='ignore')
                content = content.replace("Your Title Here", topic)
                content = content.replace("Pitch Deck", f"Презентация: {topic}")
                theme_path = self.skill_path / "assets" / "themes" / f"{theme}.css"
                if theme_path.exists():
                    theme_link = f'<link rel="stylesheet" href="{theme_path}">'
                else:
                    theme_link = f'<link rel="stylesheet" href="../../assets/themes/{theme}.css">'
                if '<head>' in content:
                    content = content.replace('<head>', f'<head>\n    {theme_link}')
                if animation:
                    anim_class = f'class="{animation}-animation"'
                    content = content.replace('<body>', f'<body {anim_class}>')
                index_path.write_text(content, encoding='utf-8')
            return index_path, None, f"HTML презентация создана"
        return None, None, "Шаблон не найден"
    
    async def convert_html_to_pptx_async(self, html_path: Path, pptx_path: Path):
        if self.apograph_available:
            from apograph import convert_html_to_pptx
            await convert_html_to_pptx(str(html_path), str(pptx_path))
            return pptx_path
        else:
            raise ImportError("apograph не установлен. Выполните: pip install apograph")
    
    def convert_html_to_pptx(self, html_path, pptx_path=None):
        html = Path(html_path)
        if not pptx_path:
            pptx_path = self.output_dir / f"{html.stem}_converted.pptx"
        else:
            pptx_path = Path(pptx_path)
        
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        result = loop.run_until_complete(self.convert_html_to_pptx_async(html, pptx_path))
        loop.close()
        return result
    
    def create_and_convert_ppt(self, topic, slides=5, theme="cyberpunk-neon", animation="fade"):
        html_path, _, msg = self.create_html_ppt(topic, slides, theme, animation)
        if html_path and html_path.exists():
            pptx_path = self.convert_html_to_pptx(html_path)
            return pptx_path, f"HTML и PPTX созданы: {pptx_path}"
        return None, msg
    
    def create_ppt(self, user_input, slides_count=5):
        topic, slides_count = self._parse_topic_and_slides(user_input, slides_count)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = self.output_dir / f"{topic[:30]}_{timestamp}.pptx"
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]
        self._generate_content_slides(prs, title_slide_layout, content_slide_layout, topic, slides_count)
        prs.save(str(output_file))
        return output_file, f"PPTX презентация создана"
    
    def _parse_topic_and_slides(self, user_input, slides_count=5):
        numbers = re.findall(r'\d+', user_input)
        if numbers:
            slides_count = int(numbers[0])
            topic = re.sub(r'\d+', '', user_input).strip()
        else:
            topic = user_input.strip()
        if not topic:
            topic = "Презентация"
        return topic, slides_count
    
    def _generate_content_slides(self, prs, title_layout, content_layout, topic, num_slides):
        slides_data = [
            {"title": f"Введение в {topic}", "content": f"Сегодня мы рассмотрим ключевые аспекты {topic}."},
            {"title": f"Основы {topic}", "content": f"• Первый важный принцип\n• Второй ключевой момент\n• Практическое применение"},
            {"title": f"Преимущества", "content": f"1. Эффективность\n2. Надёжность\n3. Масштабируемость"},
            {"title": f"Примеры использования", "content": f"Пример 1: Базовый сценарий\nПример 2: Продвинутый подход\nПример 3: Интеграция"},
            {"title": f"Заключение", "content": f"• {topic} открывает новые возможности\n• Рекомендации по внедрению\n• Дальнейшие шаги"}
        ]
        for i in range(min(num_slides, len(slides_data))):
            slide_data = slides_data[i]
            if i == 0:
                slide = prs.slides.add_slide(title_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = slide_data["title"]
                subtitle.text = slide_data["content"]
            else:
                slide = prs.slides.add_slide(content_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                title.text = slide_data["title"]
                tf = content.text_frame
                tf.text = slide_data["content"]
                for paragraph in tf.paragraphs:
                    paragraph.font.size = Pt(18)
